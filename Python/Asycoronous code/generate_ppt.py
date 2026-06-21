from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import textwrap

INPUT_PATH = "Module System in Python.txt"
OUTPUT_PPT = "Module_System_in_Python.pptx"

MAX_LINES_PER_SLIDE = 18


def add_title_slide(prs, title, subtitle=None):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_tf = slide.shapes.title
    title_tf.text = title
    if subtitle:
        try:
            subtitle_tf = slide.placeholders[1]
            subtitle_tf.text = subtitle
        except Exception:
            pass


def add_content_slide(prs, title, lines, notes=None):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(lines):
        p = body.add_paragraph() if i > 0 else body.paragraphs[0]
        p.level = 0
        run = p.add_run()
        run.text = line
        font = run.font
        font.name = "Consolas"
        font.size = Pt(14)
        font.color.rgb = RGBColor(230, 230, 230)
    if notes:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.clear()
        notes_tf.text = notes


def chunk_lines(lines, chunk_size):
    chunks = []
    cur = []
    for line in lines:
        if line.strip() == "":
            # preserve blank lines as separators
            if cur:
                chunks.append(cur)
                cur = []
            continue
        cur.append(line.rstrip())
        if len(cur) >= chunk_size:
            chunks.append(cur)
            cur = []
    if cur:
        chunks.append(cur)
    return chunks


def main():
    with open(INPUT_PATH, "r", encoding="utf-8") as f:
        raw = f.read()
    lines = raw.splitlines()

    prs = Presentation()
    prs.slide_height = Inches(7.5)
    prs.slide_width = Inches(13.33)

    # Title slide
    add_title_slide(
        prs,
        "Module System in Python",
        "Organizing, Importing, and Managing Python Code",
    )

    # Add an index slide
    intro_text = "Generated from Module System in Python.txt; contains all content, examples, diagrams, and speaker notes."
    add_content_slide(
        prs,
        "About this Presentation",
        [intro_text],
        notes="This presentation was generated from the provided module text file. All content and examples are included as text and speaker notes.",
    )

    chunks = chunk_lines(lines, MAX_LINES_PER_SLIDE)
    for idx, chunk in enumerate(chunks):
        # derive title from first line if it looks like a heading
        first = chunk[0]
        title = first if len(first) <= 60 else f"Module Content (Part {idx+1})"
        # shorten long lines for slide body but keep full in notes
        body_lines = [
            l if len(l) <= 140 else textwrap.fill(l, width=120) for l in chunk
        ]
        notes = "\n".join(chunk)
        add_content_slide(prs, title, body_lines, notes=notes)

    prs.save(OUTPUT_PPT)
    print(f"Presentation saved to {OUTPUT_PPT}")


if __name__ == "__main__":
    main()
